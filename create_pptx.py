from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT_ORANGE = RGBColor(0xFF, 0x66, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
MID_GRAY = RGBColor(0x60, 0x60, 0x60)

W = Inches(13.33)
H = Inches(7.5)
MARGIN = Inches(0.5)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullets(slide, items, left, top, width, height,
                font_size=14, color=DARK_GRAY, indent=0.25):
    """Add bulleted list items to a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = "• " + item[0]
            run.font.size = Pt(item[1])
            run.font.bold = item[2] if len(item) > 2 else False
            run.font.color.rgb = item[3] if len(item) > 3 else color
        else:
            run.text = "• " + item
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = color

def add_section_header(slide, title, subtitle=None):
    """Add a slide title + optional subtitle."""
    # Title
    add_text_box(slide, title,
                 MARGIN, Inches(0.3), W - 2*MARGIN, Inches(0.7),
                 font_size=22, bold=True, color=DARK_BLUE)
    if subtitle:
        add_text_box(slide, subtitle,
                     MARGIN, Inches(0.9), W - 2*MARGIN, Inches(0.4),
                     font_size=14, bold=False, color=ACCENT_ORANGE, italic=True)

def slide_bg(slide):
    set_bg(slide, WHITE)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Blank

# ─── SLIDE 1: Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BLUE)

# White title
add_text_box(s, "AI-First Transformation:",
             MARGIN, Inches(1.8), W - 2*MARGIN, Inches(0.9),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(s, "From Strategy to Competitive Advantage",
             MARGIN, Inches(2.7), W - 2*MARGIN, Inches(0.9),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Orange divider line
line_box = s.shapes.add_shape(1, MARGIN, Inches(3.75), Inches(2), Pt(4))
line_box.fill.solid()
line_box.fill.fore_color.rgb = ACCENT_ORANGE
line_box.line.fill.background()

add_text_box(s, "A Data-Driven Roadmap for Sustainable Growth",
             MARGIN, Inches(4.0), W - 2*MARGIN, Inches(0.5),
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(s, "Executive Review  |  April 2026",
             MARGIN, Inches(5.8), W - 2*MARGIN, Inches(0.4),
             font_size=13, bold=False, color=RGBColor(0xB0,0xC8,0xE0), align=PP_ALIGN.LEFT)

# ─── SLIDE 2: State of Industry ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "AI Adoption Is Accelerating — Is Your Organization Keeping Pace?")

# Global AI Investment
add_text_box(s, "Global AI Investment",
             MARGIN, Inches(1.4), Inches(3.8), Inches(0.35),
             font_size=15, bold=True, color=DARK_BLUE)
items1 = [
    "$320B+ invested in enterprise AI in 2025",
    "73% of Fortune 500 now have dedicated AI leadership roles",
    "AI-First companies: 40–60% efficiency gains vs. 12–18% for late adopters"
]
add_bullets(s, items1, MARGIN, Inches(1.75), Inches(3.8), Inches(1.5), font_size=13)

# Three Industries
add_text_box(s, "Three Industries Driving Growth",
             MARGIN, Inches(3.3), Inches(3.8), Inches(0.35),
             font_size=15, bold=True, color=DARK_BLUE)
items2 = [
    "Manufacturing: $78B (gas turbines, predictive maintenance, quality inspection)",
    "Financial Services: $65B (fraud detection, risk modeling, customer intelligence)",
    "Telecommunications: $52B (network optimization, AI-RAN, edge inference)"
]
add_bullets(s, items2, MARGIN, Inches(3.65), Inches(3.8), Inches(1.5), font_size=13)

# Right column — Cost of Inaction
rx = Inches(5.0)
add_text_box(s, "The Cost of Inaction",
             rx, Inches(1.4), Inches(7.8), Inches(0.35),
             font_size=15, bold=True, color=ACCENT_ORANGE)
items3 = [
    "Late adopters face 2.5× higher implementation costs by 2027",
    "Data latency gaps compound competitive disadvantage annually",
]
add_bullets(s, items3, rx, Inches(1.75), Inches(7.8), Inches(1.0), font_size=13)

# Divider
div = s.shapes.add_shape(1, rx, Inches(3.0), Inches(7.8), Pt(1))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div.line.fill.background()

add_text_box(s, "Key Insight",
             rx, Inches(3.15), Inches(7.8), Inches(0.35),
             font_size=15, bold=True, color=DARK_BLUE)
insight = "Every quarter of delay widens the gap. Leaders are pulling away — not because they started earlier, but because they invest smarter."
add_text_box(s, insight,
             rx, Inches(3.5), Inches(7.8), Inches(1.0),
             font_size=13, bold=False, color=DARK_GRAY, italic=True)

# ─── SLIDE 3: Telecom ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s,
    "Telecom: Turning Network Infrastructure into AI Factories",
    "Case Study: NVIDIA × Marvell Partnership (March 2026)")

cols = [(MARGIN, 1.4, 4.0), (Inches(5.0), 1.4, 4.0), (Inches(9.0), 1.4, 4.0)]

# Business Problem
add_text_box(s, "Business Problem",
             cols[0][0], cols[0][1], cols[0][2], Inches(0.35),
             font_size=14, bold=True, color=ACCENT_ORANGE)
add_bullets(s, [
    "Traditional RAN cannot handle AI inference workloads efficiently",
    "Networks consume 40%+ of operator energy budget",
    "AI-as-a-Service: $50B TAM by 2026 — missed revenue"
], cols[0][0], Inches(1.75), cols[0][2], Inches(1.8), font_size=12)

# AI Solution
add_text_box(s, "AI Solution",
             cols[1][0], cols[1][1], cols[1][2], Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "NVIDIA Aerial AI-RAN: AI inference at each cell tower",
    "Marvell custom XPUs + NVLink Fusion integration",
    "BlueField-3 DPU offloads networking, CPU burden ↓30%"
], cols[1][0], Inches(1.75), cols[1][2], Inches(1.8), font_size=12)

# Quantifiable Impact
add_text_box(s, "Quantifiable Impact",
             cols[2][0], cols[2][1], cols[2][2], Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "1 PetaOPs AI throughput per base station",
    "Network slice: weeks → hours (90%+ reduction)",
    "Energy efficiency: 4× improvement in TOPS/Watt",
    "NVIDIA invested $2B in Marvell"
], cols[2][0], Inches(1.75), cols[2][2], Inches(1.8), font_size=12)

# Key Lesson
div2 = s.shapes.add_shape(1, MARGIN, Inches(4.5), W - 2*MARGIN, Pt(1))
div2.fill.solid(); div2.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div2.line.fill.background()
add_text_box(s, "Key Lesson:",
             MARGIN, Inches(4.65), Inches(1.2), Inches(0.35),
             font_size=13, bold=True, color=ACCENT_ORANGE)
add_text_box(s, "Transform the network itself into an AI factory — not just the data center.",
             Inches(1.7), Inches(4.65), Inches(11.0), Inches(0.35),
             font_size=13, bold=False, color=DARK_GRAY, italic=True)

# ─── SLIDE 4: Manufacturing ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s,
    "Manufacturing: From Cost Center to Predictive Intelligence Hub",
    "Case Study: Siemens Energy Industrial AI Inspection System")

cols2 = [(MARGIN, 1.4, 4.0), (Inches(5.0), 1.4, 4.0), (Inches(9.0), 1.4, 4.0)]

add_text_box(s, "Business Problem",
             cols2[0][0], cols2[0][1], cols2[0][2], Inches(0.35),
             font_size=14, bold=True, color=ACCENT_ORANGE)
add_bullets(s, [
    "40+ manual inspection hours per gas turbine",
    "3–5% defect miss rate with human-only inspection",
    "Unplanned downtime costing $1M+ per incident"
], cols2[0][0], Inches(1.75), cols2[0][2], Inches(1.8), font_size=12)

add_text_box(s, "AI Solution",
             cols2[1][0], cols2[1][1], cols2[1][2], Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "Edge AI vision for real-time turbine blade inspection",
    "Deep learning trained on millions of annotated images",
    "NVIDIA Omniverse digital twin: training 12→6 months"
], cols2[1][0], Inches(1.75), cols2[1][2], Inches(1.8), font_size=12)

add_text_box(s, "Quantifiable Impact",
             cols2[2][0], cols2[2][1], cols2[2][2], Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "Inspection time: 40 hrs → 4–6 hrs (85% reduction)",
    "Unplanned downtime: ↓70%",
    "Defect detection: 95% → 99.6%",
    "Annual savings: ~$230M across global fleet"
], cols2[2][0], Inches(1.75), cols2[2][2], Inches(1.8), font_size=12)

div3 = s.shapes.add_shape(1, MARGIN, Inches(4.5), W - 2*MARGIN, Pt(1))
div3.fill.solid(); div3.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div3.line.fill.background()
add_text_box(s, "Key Lesson:",
             MARGIN, Inches(4.65), Inches(1.2), Inches(0.35),
             font_size=13, bold=True, color=ACCENT_ORANGE)
add_text_box(s, "Digital twins + edge AI = training speed doubled, accuracy unprecedented.",
             Inches(1.7), Inches(4.65), Inches(11.0), Inches(0.35),
             font_size=13, bold=False, color=DARK_GRAY, italic=True)

# ─── SLIDE 5: Financial Services ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s,
    "Financial Services: Real-Time AI Decisions at Global Scale",
    "Case Study: Mastercard Decision Intelligence 2.0")

cols3 = [(MARGIN, 1.4, 4.0), (Inches(5.0), 1.4, 4.0), (Inches(9.0), 1.4, 4.0)]

add_text_box(s, "Business Problem",
             cols3[0][0], cols3[0][1], cols3[0][2], Inches(0.35),
             font_size=14, bold=True, color=ACCENT_ORANGE)
add_bullets(s, [
    "$343B annual fraud loss industry-wide",
    "Rule engines: 72-hour lag detecting new attack patterns",
    "False positive rate of 1.5% damaging customer experience"
], cols3[0][0], Inches(1.75), cols3[0][2], Inches(1.8), font_size=12)

add_text_box(s, "AI Solution",
             cols3[1][0], cols3[1][1], cols3[1][2], Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "Deep neural network: 200+ real-time transaction features",
    "Transfer learning: new market protections in <48 hours",
    "GenAI summaries: 18 min → 6 min per investigation",
    "XAI module for regulatory compliance"
], cols3[1][0], Inches(1.75), cols3[1][2], Inches(1.8), font_size=12)

add_text_box(s, "Quantifiable Impact",
             cols3[2][0], cols3[2][1], cols3[2][2], Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "Fraud losses prevented: $96B/year (28% reduction)",
    "Decision latency: 300–500ms → under 50ms",
    "False positive rate: 1.5% → 0.3% (80% reduction)",
    "ROI: $1 invested avoids $17 in losses"
], cols3[2][0], Inches(1.75), cols3[2][2], Inches(1.8), font_size=12)

div4 = s.shapes.add_shape(1, MARGIN, Inches(4.5), W - 2*MARGIN, Pt(1))
div4.fill.solid(); div4.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div4.line.fill.background()
add_text_box(s, "Key Lesson:",
             MARGIN, Inches(4.65), Inches(1.2), Inches(0.35),
             font_size=13, bold=True, color=ACCENT_ORANGE)
add_text_box(s, "Real-time decisioning + XAI = regulatory approval in 6 weeks vs. 6 months.",
             Inches(1.7), Inches(4.65), Inches(11.0), Inches(0.35),
             font_size=13, bold=False, color=DARK_GRAY, italic=True)

# ─── SLIDE 6: Cross-Industry Patterns ────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "Three Patterns Separate AI Leaders from Laggards")

patterns = [
    ("Pattern 1: Edge-First Architecture",
     ACCENT_ORANGE,
     ["Process data where it lives — don't fight network latency",
      "Result: 6–10× latency improvement over cloud-dependent architectures"]),
    ("Pattern 2: Human-in-the-Loop Design",
     DARK_BLUE,
     ["Pure automation achieves 91% accuracy; human+AI achieves 98.7%",
      "AI augments experts rather than replacing them"]),
    ("Pattern 3: Platform Thinking Over Point Solutions",
     DARK_BLUE,
     ["Mastercard: Unified API for 200+ features → new rules deploy in 48 hours",
      "Reuse, don't rebuild: platform approach cuts new use case cost by 40%"]),
]

y_start = Inches(1.4)
for i, (title, color, bullets) in enumerate(patterns):
    add_text_box(s, title,
                 MARGIN, y_start, Inches(12.0), Inches(0.4),
                 font_size=16, bold=True, color=color)
    add_bullets(s, bullets, MARGIN, y_start + Inches(0.4),
                Inches(12.0), Inches(0.9), font_size=13)
    y_start += Inches(1.5)

# ─── SLIDE 7: Five Gaps ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "Close the Gap: Five Critical Gaps Blocking Your AI Transformation")

gaps = [
    ("Gap 1: Data Readiness",
     "Fragmented, inconsistent data lineage",
     "Opportunity: Establish a single source of truth that unlocks all downstream AI initiatives"),
    ("Gap 2: Infrastructure Scalability",
     "Cloud-only or legacy on-prem, not hybrid-ready",
     "Opportunity: Right-size infrastructure for each workload — sovereignty + performance"),
    ("Gap 3: Talent & Culture",
     "AI perceived as threat, not capability multiplier",
     "Opportunity: Upskill at scale; AI champions transform perception from fear to adoption"),
    ("Gap 4: Governance & Risk",
     "No model monitoring, audit trails, or XAI standards",
     "Opportunity: Build trust through explainability — unlock regulatory approval fast"),
    ("Gap 5: Use Case Prioritization",
     "Scattered initiatives, no ROI validation",
     "Opportunity: Focused portfolio of 3–5 high-ROI pilots beats 20 scattered experiments"),
]

y = Inches(1.4)
for title, problem, opp in gaps:
    add_text_box(s, title,
                 MARGIN, y, Inches(3.5), Inches(0.32),
                 font_size=13, bold=True, color=ACCENT_ORANGE)
    add_text_box(s, problem,
                 Inches(3.6), y, Inches(3.8), Inches(0.32),
                 font_size=12, bold=False, color=DARK_GRAY)
    add_text_box(s, opp,
                 Inches(7.5), y, Inches(5.3), Inches(0.32),
                 font_size=12, bold=False, color=DARK_BLUE, italic=True)
    y += Inches(0.65)

# ─── SLIDE 8: Rec 1 & 2 ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "Foundation First: Change Culture, Optimize Resources")

# Rec 1 — Culture
add_text_box(s, "Recommendation 1: Launch an AI Culture Transformation",
             MARGIN, Inches(1.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "CEO signs AI Ethics Charter",
    'Establish "AI Ambassador" program: 5–10 champions per department',
    'Reframe KPIs: replace "jobs replaced" with "collaboration efficiency"',
    "Target: 85%+ employee AI adoption within 12 months"
], MARGIN, Inches(1.8), Inches(6.0), Inches(1.8), font_size=13)

# Divider
div5 = s.shapes.add_shape(1, Inches(6.7), Inches(1.4), Pt(1), Inches(5.6))
div5.fill.solid(); div5.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div5.line.fill.background()

# Rec 2 — Budget
add_text_box(s, "Recommendation 2: Apply the 60/20/15/5 Budget Rule",
             Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE)
budget_items = [
    "60% → High-ROI pilot use cases",
    "20% → Platform infrastructure",
    "15% → Scale-up reserves",
    "5% → 前沿探索",
    "Expected: First measurable ROI within 6 months"
]
add_bullets(s, budget_items,
            Inches(7.0), Inches(1.8), Inches(6.0), Inches(1.8), font_size=13)

# Visual — 60/20/15/5 bar
bar_y = Inches(4.0)
bar_x = Inches(7.0)
bar_total = Inches(5.5)

segments = [(0.60, ACCENT_ORANGE), (0.20, DARK_BLUE), (0.15, RGBColor(0x5B,0x9B,0xC5)), (0.05, RGBColor(0xA0,0xA0,0xA0))]
cx = bar_x
for pct, color in segments:
    w = bar_total * pct
    rect = s.shapes.add_shape(1, cx, bar_y, w, Inches(0.35))
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    cx += w

labels = ["60% Pilot", "20% Platform", "15% Scale", "5% Explore"]
lx = bar_x
for pct, lbl in zip([0.60, 0.20, 0.15, 0.05], labels):
    add_text_box(s, lbl,
                 lx, bar_y + Inches(0.38), bar_total * pct, Inches(0.28),
                 font_size=11, bold=True, color=DARK_GRAY, align=PP_ALIGN.LEFT if lx == bar_x else PP_ALIGN.LEFT)
    lx += bar_total * pct

# ─── SLIDE 9: Rec 3 & 4 ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "Build the Right Foundation: Scalable Infrastructure + Data Readiness")

# Rec 3
add_text_box(s, "Recommendation 3: Deploy a Hybrid AI Architecture",
             MARGIN, Inches(1.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE)
add_bullets(s, [
    "Core: On-prem inference (data sovereignty, latency)",
    "Burst: Cloud for training and peak load",
    "MLOps pipeline with automated retraining triggers"
], MARGIN, Inches(1.8), Inches(6.0), Inches(1.4), font_size=13)

# Visual boxes for hybrid
boxes = [("On-Prem\nInference", DARK_BLUE, MARGIN, Inches(3.4)),
         ("Cloud\nTraining/Burst", RGBColor(0x5B,0x9B,0xC5), Inches(3.2), Inches(3.4)),
         ("Edge\nInference", ACCENT_ORANGE, Inches(6.0), Inches(3.4))]
for label, color, bx, by in boxes:
    r = s.shapes.add_shape(1, bx, by, Inches(2.5), Inches(0.7))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    add_text_box(s, label, bx, by + Inches(0.1), Inches(2.5), Inches(0.5),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Rec 4
add_text_box(s, "Recommendation 4: Accelerate Data Maturity",
             Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE)
levels = [
    ("L1 (0–2 mo):", "Map data lineage — quick win"),
    ("L2 (12 mo):", "Establish data ownership"),
    ("L3 (12–18 mo):", "Quality standards + MLOps baseline"),
    ("L4 (ongoing):", "Data monetization, proactive governance"),
]
ly = Inches(1.85)
for level, desc in levels:
    add_text_box(s, level, Inches(7.0), ly, Inches(1.5), Inches(0.3),
                 font_size=12, bold=True, color=ACCENT_ORANGE)
    add_text_box(s, desc, Inches(8.5), ly, Inches(4.5), Inches(0.3),
                 font_size=12, bold=False, color=DARK_GRAY)
    ly += Inches(0.42)

# Arrow progression
arrow_y = Inches(5.2)
for i, (lbl, _) in enumerate(levels):
    bx = Inches(7.0) + i * Inches(1.5)
    r = s.shapes.add_shape(1, bx, arrow_y, Inches(1.1), Inches(0.35))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT_ORANGE if i == 0 else DARK_BLUE
    r.line.fill.background()
    add_text_box(s, lbl.split(" (")[0], bx, arrow_y + Inches(0.05),
                 Inches(1.1), Inches(0.28),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE 10: Rec 5 & Roadmap ───────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "Mitigate Risks Proactively: Three-Layer Defense + 18-Month Roadmap")

# Three-Layer Defense
add_text_box(s, "Three-Layer Risk Defense",
             MARGIN, Inches(1.4), Inches(6.0), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
layers = [
    ("Layer 1 — Technical", "Model drift monitoring, A/B testing"),
    ("Layer 2 — Governance", "XAI, audit logs, fairness testing"),
    ("Layer 3 — Business Continuity", "Human-in-loop, degradation plans"),
]
ly2 = Inches(1.75)
for title, desc in layers:
    add_text_box(s, title + ":",
                 MARGIN, ly2, Inches(2.0), Inches(0.3),
                 font_size=12, bold=True, color=ACCENT_ORANGE)
    add_text_box(s, desc,
                 MARGIN + Inches(2.0), ly2, Inches(4.0), Inches(0.3),
                 font_size=12, bold=False, color=DARK_GRAY)
    ly2 += Inches(0.38)

# Divider
div6 = s.shapes.add_shape(1, Inches(6.7), Inches(1.4), Pt(1), Inches(5.6))
div6.fill.solid(); div6.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div6.line.fill.background()

# 18-Month Roadmap
add_text_box(s, "18-Month Roadmap",
             Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
roadmap = [
    ("Q2 2026", "Data assessment + AI Ethics Charter + Use case selection"),
    ("Q3 2026", "MLOps platform + First pilot go-live"),
    ("Q4 2026", "Scale to 3–5 use cases, First ROI report"),
    ("Q1 2027", "Hybrid infrastructure complete"),
    ("Q2 2027", "Organization-wide adoption >80%"),
]
ry = Inches(1.8)
for qtr, desc in roadmap:
    # Quarter badge
    r = s.shapes.add_shape(1, Inches(7.0), ry, Inches(1.1), Inches(0.32))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT_ORANGE
    r.line.fill.background()
    add_text_box(s, qtr, Inches(7.0), ry + Inches(0.05),
                 Inches(1.1), Inches(0.28),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, desc, Inches(8.2), ry + Inches(0.03),
                 Inches(4.6), Inches(0.3),
                 font_size=12, bold=False, color=DARK_GRAY)
    ry += Inches(0.48)

# ─── SLIDE 11: ROI ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "The ROI Case: Every Dollar Invested Returns 12–17× Within 36 Months")

# Financial Impact table-style
years = [
    ("Year 1", "Efficiency Gains 15–20%\nCost Reduction $X M", LIGHT_GRAY),
    ("Year 2", "Efficiency Gains 35–45%\nCost Reduction $2.5X M", WHITE),
    ("Year 3", "Efficiency Gains 50–60%\nCost Reduction $4X M\nNew Revenue $3X M", LIGHT_GRAY),
]
yw = Inches(1.4)
for yr, text, bg in years:
    r = s.shapes.add_shape(1, yw, Inches(1.6), Inches(3.8), Inches(1.4))
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)
    add_text_box(s, yr, yw + Inches(0.1), Inches(1.65),
                 Inches(3.6), Inches(0.35),
                 font_size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s, text, yw + Inches(0.1), Inches(2.0),
                 Inches(3.6), Inches(0.9),
                 font_size=13, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    yw += Inches(4.0)

# Resource Requirements
add_text_box(s, "Resource Requirements",
             MARGIN, Inches(3.3), Inches(12.0), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
reqs = [
    "Core team: 8–12 FTEs (data scientists, MLOps engineers, AI architects)",
    "Annual AI budget: 1.5–2% of revenue",
    "Executive sponsorship required",
    "Timeline to first value: 6–9 months"
]
add_bullets(s, reqs, MARGIN, Inches(3.65), Inches(12.0), Inches(1.6), font_size=13)

# ─── SLIDE 12: Next Steps ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_section_header(s, "Your Decision: Pilot in Q2 or Risk Falling Behind")

# Immediate Actions
add_text_box(s, "Immediate Actions (Next 30 Days)",
             MARGIN, Inches(1.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE)
actions = [
    "Schedule AI readiness assessment workshop (Week 1–2)",
    "Identify 3 candidate pilot use cases (Week 2–3)",
    "Secure executive sponsorship and budget allocation (Week 3–4)",
    "Form AI Transformation Steering Committee (Week 4)",
]
add_bullets(s, actions, MARGIN, Inches(1.8), Inches(6.0), Inches(1.8), font_size=13)

# Key Milestone
div7 = s.shapes.add_shape(1, Inches(6.7), Inches(1.4), Pt(1), Inches(5.6))
div7.fill.solid(); div7.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div7.line.fill.background()

add_text_box(s, "Key Milestones",
             Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=ACCENT_ORANGE)
milestones = [
    "First pilot live within 6 months",
    "Measurable ROI within 12 months",
    "Organization-wide adoption >80% by Q2 2027"
]
add_bullets(s, milestones, Inches(7.0), Inches(1.8), Inches(6.0), Inches(1.2), font_size=13)

# CTA Banner
cta = s
# CTA Banner at bottom
cta_rect = s.shapes.add_shape(1, MARGIN, Inches(4.5), W - 2*MARGIN, Inches(1.1))
cta_rect.fill.solid(); cta_rect.fill.fore_color.rgb = DARK_BLUE
cta_rect.line.fill.background()
add_text_box(s, "The time to act is now. Leaders who pilot in Q2 2026 will set the pace for the decade.",
             MARGIN + Inches(0.2), Inches(4.6), W - 2*MARGIN - Inches(0.4), Inches(0.9),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s, "Q&A",
             MARGIN, Inches(5.9), W - 2*MARGIN, Inches(0.5),
             font_size=20, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# ─── SAVE ─────────────────────────────────────────────────────────────────────
out = "/root/.openclaw/workspace/zhugeliang/AI_First_Transformation_Strategy_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
