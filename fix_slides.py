#!/usr/bin/env python3
"""Fix slides 3, 4, 5 of AI_First_Transformation_Strategy_2026.pptx"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
# Using shape type 1 for rectangle per spec
from pptx.oxml.ns import qn
from lxml import etree
import copy

# Color constants
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xFF, 0x66, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Layout constants
LEFT_MARGIN = Emu(457200)
COL_WIDTH = Emu(2819400)
COL_GAP = Emu(548640)
COL3_WIDTH = Emu(2336760)
COL1_X = LEFT_MARGIN
COL2_X = Emu(3505200)
COL3_X = Emu(6553200)

SLIDE_WIDTH = Emu(9144000)

def clear_slide(slide):
    """Delete all shapes from slide."""
    # Get all shape IDs to delete
    shapes_to_delete = list(slide.shapes)
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

def add_header_shape(slide, text, left, top, width, height):
    """Add a filled rectangle header with centered white text."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    # Clear default paragraph
    tf.paragraphs[0].text = ""
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Vertically center
    tf.anchor = 1  # middle

    return shape

def add_textbox(slide, text, left, top, width, height, font_size, bold=False,
                italic=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Add a textbox with specified formatting."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tb

def add_bullet_textbox(slide, bullets, left, top, width, height, font_size=12):
    """Add a textbox with bullet points."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return tb

def add_orange_line(slide, left, top, width, height):
    """Add a thin orange rectangle as a divider line."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    return shape

def add_orange_accent(slide, left, top, width, height):
    """Add orange accent rectangle."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    return shape

def add_key_lesson_box(slide, lesson_text, left, top, width, height):
    """Add key lesson box with bold label and normal text."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    # First paragraph: "Key Lesson:"
    p1 = tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "Key Lesson: "
    run1.font.bold = True
    run1.font.size = Pt(12)
    run1.font.color.rgb = DARK_BLUE

    # Second paragraph: the actual lesson
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = lesson_text
    run2.font.size = Pt(12)
    run2.font.color.rgb = DARK_GRAY

    return tb

def build_slide_3(slide):
    """Slide 3: Telecom NVIDIA x Marvell AI-RAN"""
    # Title
    add_textbox(slide,
               "Telecom: Turning Network Infrastructure into AI Factories",
               Emu(457200), Emu(274320), Emu(11277600), Emu(640080),
               font_size=20, bold=True, color=DARK_BLUE)

    # Subtitle
    add_textbox(slide,
               "Case Study: NVIDIA × Marvell Partnership (March 2026)",
               Emu(457200), Emu(822960), Emu(11277600), Emu(365760),
               font_size=14, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    # Three column headers
    add_header_shape(slide, "Business Problem", COL1_X, Emu(914400), COL_WIDTH, Emu(457200))
    add_header_shape(slide, "AI Solution", COL2_X, Emu(914400), COL_WIDTH, Emu(457200))
    add_header_shape(slide, "Quantifiable Impact", COL3_X, Emu(914400), COL3_WIDTH, Emu(457200))

    # Content 1: Business Problem
    add_bullet_textbox(slide,
        ["Traditional RAN architecture can't handle AI inference workloads",
         "Networks consume 40%+ of operator energy budget",
         "AI-as-a-Service: $50B TAM by 2026 — missed revenue opportunity"],
        COL1_X, Emu(1447800), COL_WIDTH, Emu(2438400), font_size=12)

    # Content 2: AI Solution
    add_bullet_textbox(slide,
        ["NVIDIA Aerial AI-RAN: AI inference at each cell tower",
         "Marvell custom XPUs + NVLink Fusion integration",
         "BlueField-3 DPU offloads networking, CPU burden ↓30%"],
        COL2_X, Emu(1447800), COL_WIDTH, Emu(2438400), font_size=12)

    # Content 3: Quantifiable Impact
    add_bullet_textbox(slide,
        ["1 PetaOPs AI throughput per base station",
         "Network slice: weeks → hours (90%+ reduction)",
         "Energy efficiency: 4× improvement in TOPS/Watt",
         "NVIDIA invested $2B in Marvell"],
        COL3_X, Emu(1447800), COL3_WIDTH, Emu(2438400), font_size=12)

    # Orange accent (top right)
    add_orange_accent(slide, Emu(8500000), Emu(914400), Emu(457200), Emu(91440))

    # Orange divider line above key lesson
    add_orange_line(slide, Emu(457200), Emu(4038600), Emu(11277600), Emu(15240))

    # Key Lesson box
    add_key_lesson_box(slide,
        "Transform the network itself into an AI factory — not just the data center.",
        Emu(457200), Emu(4114800), Emu(11277600), Emu(548640))

def build_slide_4(slide):
    """Slide 4: Manufacturing Siemens Energy"""
    # Title
    add_textbox(slide,
               "Manufacturing: From Cost Center to Predictive Intelligence Hub",
               Emu(457200), Emu(274320), Emu(11277600), Emu(640080),
               font_size=20, bold=True, color=DARK_BLUE)

    # Subtitle
    add_textbox(slide,
               "Case Study: Siemens Energy Industrial AI Inspection System",
               Emu(457200), Emu(822960), Emu(11277600), Emu(365760),
               font_size=14, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    # Three column headers
    add_header_shape(slide, "Business Problem", COL1_X, Emu(914400), COL_WIDTH, Emu(457200))
    add_header_shape(slide, "AI Solution", COL2_X, Emu(914400), COL_WIDTH, Emu(457200))
    add_header_shape(slide, "Quantifiable Impact", COL3_X, Emu(914400), COL3_WIDTH, Emu(457200))

    # Content 1: Business Problem
    add_bullet_textbox(slide,
        ["40+ manual inspection hours per gas turbine",
         "3–5% defect miss rate with human-only inspection",
         "Unplanned downtime costing $1M+ per incident"],
        COL1_X, Emu(1447800), COL_WIDTH, Emu(2438400), font_size=12)

    # Content 2: AI Solution
    add_bullet_textbox(slide,
        ["Edge AI vision for real-time turbine blade inspection",
         "Deep learning model trained on millions of annotated images",
         "NVIDIA Omniverse digital twin: training time 12→6 months"],
        COL2_X, Emu(1447800), COL_WIDTH, Emu(2438400), font_size=12)

    # Content 3: Quantifiable Impact
    add_bullet_textbox(slide,
        ["Inspection time: 40 hrs → 4–6 hrs (85% reduction)",
         "Unplanned downtime: ↓70%",
         "Defect detection rate: 95% → 99.6%",
         "Annual maintenance savings: ~$230M across global fleet"],
        COL3_X, Emu(1447800), COL3_WIDTH, Emu(2438400), font_size=12)

    # Orange accent
    add_orange_accent(slide, Emu(8500000), Emu(914400), Emu(457200), Emu(91440))

    # Orange divider line
    add_orange_line(slide, Emu(457200), Emu(4038600), Emu(11277600), Emu(15240))

    # Key Lesson box
    add_key_lesson_box(slide,
        "Digital twins + edge AI = training speed doubled, accuracy unprecedented.",
        Emu(457200), Emu(4114800), Emu(11277600), Emu(548640))

def build_slide_5(slide):
    """Slide 5: Financial Services Mastercard"""
    # Title
    add_textbox(slide,
               "Financial Services: Real-Time AI Decisions at Global Scale",
               Emu(457200), Emu(274320), Emu(11277600), Emu(640080),
               font_size=20, bold=True, color=DARK_BLUE)

    # Subtitle
    add_textbox(slide,
               "Case Study: Mastercard Decision Intelligence 2.0",
               Emu(457200), Emu(822960), Emu(11277600), Emu(365760),
               font_size=14, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    # Three column headers
    add_header_shape(slide, "Business Problem", COL1_X, Emu(914400), COL_WIDTH, Emu(457200))
    add_header_shape(slide, "AI Solution", COL2_X, Emu(914400), COL_WIDTH, Emu(457200))
    add_header_shape(slide, "Quantifiable Impact", COL3_X, Emu(914400), COL3_WIDTH, Emu(457200))

    # Content 1: Business Problem
    add_bullet_textbox(slide,
        ["$343B annual fraud loss industry-wide",
         "Rule engines: 72-hour lag detecting new attack patterns",
         "False positive rate of 1.5% damaging customer experience"],
        COL1_X, Emu(1447800), COL_WIDTH, Emu(2438400), font_size=12)

    # Content 2: AI Solution
    add_bullet_textbox(slide,
        ["Deep neural network: 200+ real-time transaction features",
         "Transfer learning: new market protections in <48 hours",
         "GenAI summaries: 18 min → 6 min per investigation",
         "XAI module for regulatory compliance"],
        COL2_X, Emu(1447800), COL_WIDTH, Emu(2438400), font_size=12)

    # Content 3: Quantifiable Impact
    add_bullet_textbox(slide,
        ["Fraud losses prevented: $96B/year (28% reduction)",
         "Decision latency: 300–500ms → under 50ms",
         "False positive rate: 1.5% → 0.3% (80% reduction)",
         "ROI: $1 invested avoids $17 in losses"],
        COL3_X, Emu(1447800), COL3_WIDTH, Emu(2438400), font_size=12)

    # Orange accent
    add_orange_accent(slide, Emu(8500000), Emu(914400), Emu(457200), Emu(91440))

    # Orange divider line
    add_orange_line(slide, Emu(457200), Emu(4038600), Emu(11277600), Emu(15240))

    # Key Lesson box
    add_key_lesson_box(slide,
        "Real-time decisioning + XAI = regulatory approval in 6 weeks vs. 6 months.",
        Emu(457200), Emu(4114800), Emu(11277600), Emu(548640))

def main():
    pptx_path = "/root/.openclaw/workspace/zhugeliang/AI_First_Transformation_Strategy_2026.pptx"
    prs = Presentation(pptx_path)

    print(f"Total slides: {len(prs.slides)}")

    # Fix slides 3, 4, 5 (index 2, 3, 4)
    for idx in [2, 3, 4]:
        slide = prs.slides[idx]
        print(f"Clearing and rebuilding slide {idx + 1}...")
        clear_slide(slide)

    # Rebuild slide 3
    print("Building slide 3 (Telecom)...")
    build_slide_3(prs.slides[2])

    # Rebuild slide 4
    print("Building slide 4 (Manufacturing)...")
    build_slide_4(prs.slides[3])

    # Rebuild slide 5
    print("Building slide 5 (Financial Services)...")
    build_slide_5(prs.slides[4])

    # Save
    prs.save(pptx_path)

    import os
    size = os.path.getsize(pptx_path)
    print(f"\nFile saved: {pptx_path}")
    print(f"File size: {size} bytes")

if __name__ == "__main__":
    main()
